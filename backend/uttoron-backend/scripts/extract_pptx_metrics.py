#!/usr/bin/env python3
"""Extract four objective, deterministic metrics from a .pptx file.

Usage: python3 extract_pptx_metrics.py <path-to-pptx>

Always prints exactly one JSON object to stdout and exits 0 (success or
"failed" parse status both count as a completed run — a Node caller just
needs to JSON.parse stdout, never inspect the exit code). Metric
definitions match the implementation brief exactly:

  avg_words_per_slide   - mean, across slides, of summed word counts of
                           every text frame on that slide
  avg_bullets_per_slide - mean, across slides, of non-empty paragraphs
                           across every text frame on that slide
  image_to_text_ratio   - fraction of slides containing >=1 picture shape
  distinct_font_count   - count of distinct non-null run.font.name values
                           across every run in the deck

No ML, no scoring, no judgment — counting and arithmetic only.
"""
import json
import sys

def extract(path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(path)
    slides = list(prs.slides)
    slide_count = len(slides)
    if slide_count == 0:
        raise ValueError("empty deck (0 slides)")

    total_words = 0
    total_bullets = 0
    slides_with_picture = 0
    fonts = set()

    for slide in slides:
        slide_has_picture = False
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_has_picture = True
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = shape.text_frame
            total_words += len(tf.text.split())
            for paragraph in tf.paragraphs:
                if paragraph.text.strip() != "":
                    total_bullets += 1
                for run in paragraph.runs:
                    name = run.font.name
                    if name:
                        fonts.add(name)
        if slide_has_picture:
            slides_with_picture += 1

    return {
        "parse_status": "success",
        "avg_words_per_slide": total_words / slide_count,
        "avg_bullets_per_slide": total_bullets / slide_count,
        "image_to_text_ratio": slides_with_picture / slide_count,
        "distinct_font_count": len(fonts),
        "slide_count": slide_count,
    }

def main():
    if len(sys.argv) != 2:
        print(json.dumps({"parse_status": "failed", "error": "usage: extract_pptx_metrics.py <path>"}))
        return
    try:
        result = extract(sys.argv[1])
    except Exception as e:
        result = {"parse_status": "failed", "error": str(e)}
    print(json.dumps(result))

if __name__ == "__main__":
    main()
